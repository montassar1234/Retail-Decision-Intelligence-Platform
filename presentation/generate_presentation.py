from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
PROCESSED_DIR = PROJECT_ROOT / "data" / "processed"
OUTPUT_PATH = PROJECT_ROOT / "presentation" / "Smart_Retail_BI_AI_Presentation.pptx"

COLORS = {
    "navy": RGBColor(18, 42, 66),
    "teal": RGBColor(33, 128, 141),
    "gold": RGBColor(230, 184, 76),
    "light": RGBColor(240, 244, 248),
    "dark": RGBColor(35, 43, 51),
    "white": RGBColor(255, 255, 255),
    "red": RGBColor(191, 64, 64),
}


def load_inputs() -> dict[str, object]:
    return {
        "kpis": json.loads((PROCESSED_DIR / "kpis.json").read_text(encoding="utf-8")),
        "metrics": json.loads((PROCESSED_DIR / "model_metrics.json").read_text(encoding="utf-8")),
        "monthly": pd.read_csv(PROCESSED_DIR / "monthly_performance.csv"),
        "category": pd.read_csv(PROCESSED_DIR / "category_performance.csv"),
        "country": pd.read_csv(PROCESSED_DIR / "country_performance.csv"),
        "future": pd.read_csv(PROCESSED_DIR / "future_forecast.csv"),
        "anomalies": pd.read_csv(PROCESSED_DIR / "anomalies.csv"),
        "segments": pd.read_csv(PROCESSED_DIR / "product_segments.csv"),
        "recommendations": pd.read_csv(PROCESSED_DIR / "recommendations.csv"),
    }


def style_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    para = title_box.text_frame.paragraphs[0]
    para.text = title
    para.font.size = Pt(28)
    para.font.bold = True
    para.font.color.rgb = COLORS["navy"]
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12), Inches(0.5))
        ps = sub_box.text_frame.paragraphs[0]
        ps.text = subtitle
        ps.font.size = Pt(12)
        ps.font.color.rgb = COLORS["teal"]


def add_bullets(slide, items: list[str], left: float = 0.8, top: float = 1.6, width: float = 5.8, height: float = 4.8) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20 if idx == 0 and len(items) <= 3 else 18)
        p.font.color.rgb = COLORS["dark"]
        p.space_after = Pt(10)


def add_kpi_card(slide, left: float, top: float, title: str, value: str, fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.8), Inches(1.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    text_frame = shape.text_frame
    text_frame.clear()
    p1 = text_frame.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = COLORS["white"]
    p1.alignment = PP_ALIGN.CENTER
    p2 = text_frame.add_paragraph()
    p2.text = value
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = COLORS["white"]
    p2.alignment = PP_ALIGN.CENTER


def add_chart(slide, chart_type, categories, series_name, values, left, top, width, height):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    chart = slide.shapes.add_chart(chart_type, Inches(left), Inches(top), Inches(width), Inches(height), chart_data).chart
    chart.has_legend = chart_type == XL_CHART_TYPE.PIE
    if chart_type != XL_CHART_TYPE.PIE:
        chart.value_axis.has_major_gridlines = True
        chart.category_axis.tick_labels.font.size = Pt(10)
        chart.value_axis.tick_labels.font.size = Pt(10)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = COLORS["teal"]
    return chart


def build_presentation(data: dict[str, object]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    kpis = data["kpis"]
    metrics = data["metrics"]
    monthly = data["monthly"]
    category = data["category"]
    country = data["country"]
    future = data["future"]
    anomalies = data["anomalies"]
    segments = data["segments"]
    recommendations = data["recommendations"]

    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["navy"]
    bg.line.color.rgb = COLORS["navy"]
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.5), Inches(1.2))
    para = title_box.text_frame.paragraphs[0]
    para.text = "Smart Retail Demand & Decision Intelligence"
    para.font.size = Pt(28)
    para.font.bold = True
    para.font.color.rgb = COLORS["white"]
    sub = title_box.text_frame.add_paragraph()
    sub.text = "AI + BI Academic Project Presentation"
    sub.font.size = Pt(18)
    sub.font.color.rgb = COLORS["gold"]
    details = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11), Inches(1))
    details_para = details.text_frame.paragraphs[0]
    details_para.text = f"Dataset period: {kpis['date_range_start']} to {kpis['date_range_end']} | Open-source Online Retail dataset"
    details_para.font.size = Pt(16)
    details_para.font.color.rgb = COLORS["light"]

    slide = prs.slides.add_slide(blank)
    style_title(slide, "1. Context and Problem")
    add_bullets(
        slide,
        [
            "Retail managers often use static reports that describe the past but do not predict the future.",
            "This makes inventory, pricing, and operational decisions reactive instead of proactive.",
            "The project combines BI dashboards and AI models to support smarter retail decisions.",
        ],
        width=6.2,
    )
    add_kpi_card(slide, 7.1, 1.9, "Orders Analysed", f"{kpis['orders_count']:,}", COLORS["teal"])
    add_kpi_card(slide, 10.0, 1.9, "Countries", f"{kpis['countries_count']}", COLORS["gold"])
    add_kpi_card(slide, 7.1, 3.5, "Customers", f"{kpis['customers_count']:,}", COLORS["navy"])
    add_kpi_card(slide, 10.0, 3.5, "30-Day Growth", f"{kpis['sales_growth_pct']:.1f}%", COLORS["red"])

    slide = prs.slides.add_slide(blank)
    style_title(slide, "2. Proposed Solution")
    add_bullets(
        slide,
        [
            "BI layer: executive KPIs, category analysis, country performance, and sales trends.",
            "AI layer: demand forecasting, anomaly detection, product segmentation, and recommendations.",
            "Output: an interactive decision-support dashboard for managers and analysts.",
        ],
        width=6.2,
    )
    add_kpi_card(slide, 7.0, 1.8, "Total Revenue", f"${kpis['total_revenue']/1_000_000:.2f}M", COLORS["teal"])
    add_kpi_card(slide, 9.9, 1.8, "Total Profit", f"${kpis['total_profit']/1_000_000:.2f}M", COLORS["navy"])
    add_kpi_card(slide, 7.0, 3.4, "Avg Order Value", f"${kpis['average_order_value']:.0f}", COLORS["gold"])
    add_kpi_card(slide, 9.9, 3.4, "Forecast Next 30 Days", f"${kpis['forecast_30d_revenue']/1_000_000:.2f}M", COLORS["red"])

    slide = prs.slides.add_slide(blank)
    style_title(slide, "3. Dataset")
    add_bullets(
        slide,
        [
            "Open-source Online Retail transaction dataset.",
            "Main fields: invoice number, product, quantity, date, unit price, customer, and country.",
            "Cleaned by removing duplicates and invalid transactions, then enriched with revenue, profit, region, and category fields.",
        ],
        width=5.9,
    )
    top_countries = country.head(5)
    add_chart(
        slide,
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        top_countries["Country"].tolist(),
        "Revenue",
        top_countries["revenue"].tolist(),
        6.8,
        1.8,
        5.8,
        4.2,
    )

    slide = prs.slides.add_slide(blank)
    style_title(slide, "4. Technical Workflow")
    add_bullets(
        slide,
        [
            "1. Data cleaning and preprocessing",
            "2. KPI and business analysis",
            "3. Forecasting model",
            "4. Anomaly detection",
            "5. Product segmentation",
            "6. Dashboard development",
        ],
        width=5.2,
    )
    add_chart(
        slide,
        XL_CHART_TYPE.LINE,
        monthly["Month"].tolist(),
        "Monthly Revenue",
        monthly["revenue"].tolist(),
        6.0,
        1.7,
        6.5,
        4.3,
    )

    slide = prs.slides.add_slide(blank)
    style_title(slide, "5. BI Layer Results")
    add_kpi_card(slide, 0.8, 1.5, "Revenue", f"${kpis['total_revenue']/1_000_000:.2f}M", COLORS["teal"])
    add_kpi_card(slide, 3.8, 1.5, "Profit", f"${kpis['total_profit']/1_000_000:.2f}M", COLORS["navy"])
    add_kpi_card(slide, 6.8, 1.5, "Orders", f"{kpis['orders_count']:,}", COLORS["gold"])
    add_kpi_card(slide, 9.8, 1.5, "AOV", f"${kpis['average_order_value']:.0f}", COLORS["red"])
    top_categories = category.head(5)
    add_chart(
        slide,
        XL_CHART_TYPE.BAR_CLUSTERED,
        top_categories["Category"].tolist(),
        "Revenue",
        top_categories["revenue"].tolist(),
        0.9,
        3.0,
        5.7,
        3.6,
    )
    add_bullets(
        slide,
        [
            "BI dashboards highlight the most valuable categories and strongest countries.",
            "The same processed data can also feed a Power BI dashboard if required for class.",
        ],
        left=7.0,
        top=3.2,
        width=5.5,
        height=2.4,
    )

    slide = prs.slides.add_slide(blank)
    style_title(slide, "6. AI Layer Results")
    add_bullets(
        slide,
        [
            f"Forecasting MAE: {metrics['mae']:.0f}",
            f"Forecasting RMSE: {metrics['rmse']:.0f}",
            f"Forecasting MAPE: {metrics['mape']:.1f}%",
            f"Next 30-day revenue forecast: ${kpis['forecast_30d_revenue']:,.0f}",
            f"Anomalies detected: {len(anomalies)}",
            f"Product segments created: {segments['segment_name'].nunique()}",
        ],
        width=5.1,
    )
    future_short = future.head(10).copy()
    future_short["date"] = pd.to_datetime(future_short["date"]).dt.strftime("%m-%d")
    add_chart(
        slide,
        XL_CHART_TYPE.LINE,
        future_short["date"].tolist(),
        "Predicted Revenue",
        future_short["predicted_revenue"].tolist(),
        6.0,
        1.8,
        6.4,
        4.2,
    )

    slide = prs.slides.add_slide(blank)
    style_title(slide, "7. Product Segmentation and Anomalies")
    segment_counts = segments["segment_name"].value_counts()
    add_chart(
        slide,
        XL_CHART_TYPE.PIE,
        segment_counts.index.tolist(),
        "Products",
        segment_counts.values.tolist(),
        0.8,
        1.7,
        5.3,
        4.4,
    )
    anomaly_types = anomalies["anomaly_type"].value_counts()
    add_chart(
        slide,
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        anomaly_types.index.tolist(),
        "Count",
        anomaly_types.values.tolist(),
        6.8,
        1.8,
        5.5,
        3.2,
    )
    add_bullets(
        slide,
        [
            "Segmentation groups products into high performers, seasonal opportunities, and at-risk items.",
            "Anomaly detection flags abnormal spikes and drops that may indicate promotions, stock issues, or reporting errors.",
        ],
        left=6.8,
        top=5.2,
        width=5.6,
        height=1.5,
    )

    slide = prs.slides.add_slide(blank)
    style_title(slide, "8. Business Recommendations")
    rec_items = [
        f"{row['theme']}: {row['recommendation']} ({row['evidence']})"
        for _, row in recommendations.iterrows()
    ]
    add_bullets(slide, rec_items, width=11.5, height=4.8)

    slide = prs.slides.add_slide(blank)
    style_title(slide, "9. Deliverables")
    add_bullets(
        slide,
        [
            "Runnable Python pipeline",
            "Interactive Streamlit BI dashboard",
            "Processed CSV outputs for reporting or Power BI",
            "Forecast, anomaly, and segmentation outputs",
            "Final report and project documentation",
            "GitHub-ready repository structure",
        ],
        width=5.8,
    )
    add_bullets(
        slide,
        [
            "Target users:",
            "Sales manager",
            "Operations manager",
            "Business analyst",
            "Supply planner",
        ],
        left=7.1,
        top=1.8,
        width=4.5,
        height=3.6,
    )

    slide = prs.slides.add_slide(blank)
    style_title(slide, "10. Conclusion")
    add_bullets(
        slide,
        [
            "The project combines descriptive BI and predictive AI in one realistic retail solution.",
            "It moves business users from historical reporting to forward-looking decision support.",
            "The approach is feasible, explainable, and suitable for academic presentation and portfolio use.",
        ],
        width=11.2,
        height=3.5,
    )
    closing = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.8))
    cp = closing.text_frame.paragraphs[0]
    cp.text = "Thank you"
    cp.font.size = Pt(26)
    cp.font.bold = True
    cp.font.color.rgb = COLORS["teal"]
    cp.alignment = PP_ALIGN.CENTER

    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    build_presentation(load_inputs())
    print(f"Presentation created: {OUTPUT_PATH}")
